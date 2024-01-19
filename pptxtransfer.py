import os
import sys
import pptx
import logging
import tempfile
from PIL import Image
from pptx.util import Inches, Pt
from gtts import gTTS
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

def print_instructions():
    print("Instructions for using pptxtransfer.py:")
    print("- Ensure all dependencies (python-pptx, Pillow, gTTS, moviepy) are installed.")
    print("- Run the script with python pptxtransfer.py. Follow the prompts to input the PowerPoint file path.")
    print("- The script will convert the PowerPoint into a video and save it to the specified output path.")

def check_dependencies():
    dependencies = {"pptx": "python-pptx", "PIL": "Pillow", "gtts": "gTTS", "moviepy": "moviepy"}
    missing_dependencies = []
    for import_name, package_name in dependencies.items():
        try:
            __import__(import_name)
        except ImportError:
            missing_dependencies.append(package_name)
    if missing_dependencies:
        print("The following dependencies are missing:")
        for dep in missing_dependencies:
            print(f"- {dep}")
        print("\nPlease install them using the following command:")
        print(f"pip install {' '.join(missing_dependencies)}")
        return False
    return True

def validate_file_path(file_path, expected_extension, check_exists=True):
    if check_exists and not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")
    if not file_path.endswith(expected_extension):
        raise ValueError(f"The file {file_path} does not have the expected {expected_extension} extension.")
    return True

def export_slide_as_image(slide, idx, presentation):
    img_path = tempfile.mktemp(suffix=f"_{idx}.png")
    slide_width, slide_height = presentation.slide_width, presentation.slide_height
    img = Image.new("RGB", (int(slide_width.pt), int(slide_height.pt)), "white")
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue  # Only export non-placeholder shapes
        left, top, width, height = shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt
        # Add more logic here if needed to handle different shape types
    img.save(img_path)
    return img_path

def check_for_speaker_notes(presentation):
    for slide in presentation.slides:
        if slide.notes_slide and slide.notes_slide.notes_text_frame.text:
            return True
    return False

def get_user_input_for_skipping_slides():
    user_input = input("No slides with speaker notes found. Do you want to display all slides for a specified duration? (yes/no): ").strip().lower()
    if user_input == "yes":
        while True:
            try:
                slide_duration = int(input("Enter the number of seconds to display each slide: "))
                break
            except ValueError:
                print("Please enter a valid integer for the number of seconds.")
        return False, slide_duration
    return True, 0

def pptx_to_video(pptx_path, output_path):
    logging.info("Starting the conversion process.")
    try:
        validate_file_path(pptx_path, '.pptx')
        validate_file_path(output_path, '.mp4', check_exists=False)
    except (FileNotFoundError, ValueError) as e:
        logging.error(f"File path validation error: {e}")
        sys.exit(1)
    if not check_dependencies():
        logging.error("Missing dependencies.")
        sys.exit(1)
    presentation = pptx.Presentation(pptx_path)
    has_speaker_notes = check_for_speaker_notes(presentation)
    skip_slides, slide_duration = (False, 5)
    if not has_speaker_notes:
        skip_slides, slide_duration = get_user_input_for_skipping_slides()
    temp_audio_files, temp_image_files = [], []
    for i, slide in enumerate(presentation.slides):
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''
            if notes:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as audio_file:
                    gTTS(notes).save(audio_file.name)
                    temp_audio_files.append(audio_file.name)
            image_file_name = export_slide_as_image(slide, i, presentation)
            temp_image_files.append(image_file_name)
        except Exception as e:
            logging.error(f"Error processing slide {i}: {e}")
    video_clips = []
    for i, image_filename in enumerate(temp_image_files):
        try:
            img_clip = ImageClip(image_filename)
            if i < len(temp_audio_files):
                audio_clip = AudioFileClip(temp_audio_files[i])
                img_clip = img_clip.set_duration(audio_clip.duration).set_audio(audio_clip)
            else:
                img_clip = img_clip.set_duration(slide_duration)
            video_clips.append(img_clip)
        except Exception as e:
            logging.error(f"Error creating video clip for slide {i}: {e}")
    try:
        final_clip = concatenate_videoclips(video_clips)
        final_clip.write_videofile(output_path, fps=24)
        logging.info(f"Output video saved to {output_path}")
    except Exception as e:
        logging.error(f"Error concatenating video clips: {e}")
    for filename in temp_audio_files + temp_image_files:
        if os.path.exists(filename):
            os.remove(filename)

if __name__ == "__main__":
    pptx_to_video('example.pptx', 'output_video.mp4')
